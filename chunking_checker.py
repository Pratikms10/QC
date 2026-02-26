from pptx import Presentation
from pptx.dml.color import RGBColor
from sentence_transformers import SentenceTransformer, util
import pandas as pd
import re

model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

def get_slide_text_with_position(slide):
    shape_data = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        top = shape.top or 0
        left = shape.left or 0
        for para in shape.text_frame.paragraphs:
            run_text = ''.join([run.text for run in para.runs]).strip()
            font_rgb = None
            if para.runs and para.runs[0].font.color and para.runs[0].font.color.type == 1:
                font_rgb = para.runs[0].font.color.rgb
            if font_rgb != RGBColor(242, 103, 34) and run_text:
                shape_data.append((top, left, run_text))
    # Sort by top and left positions to get visual order
    shape_data.sort(key=lambda x: (x[0], x[1]))
    return [text for _, _, text in shape_data]

def get_vo_text(slide):
    if not slide.has_notes_slide:
        return []
    notes_text = slide.notes_slide.notes_text_frame.text
    vo_section = re.search(r'VO:(.*?)(Image Link:|Instructions to GD:|$)', notes_text, re.DOTALL | re.IGNORECASE)
    if not vo_section:
        return []
    vo_clean = vo_section.group(1)
    return [line.strip("-• \n") for line in vo_clean.split("\n") if line.strip()]

def compare_point_to_vo(point, vo_lines):
    if not vo_lines:
        return ("", 0.0, "Missing", "No VO content")
    point_emb = model.encode(point, convert_to_tensor=True)
    vo_embs = model.encode(vo_lines, convert_to_tensor=True)
    cosine_scores = util.cos_sim(point_emb, vo_embs)[0]
    best_score = float(cosine_scores.max())
    best_idx = int(cosine_scores.argmax())
    best_match = vo_lines[best_idx]
    if point.strip().lower() == best_match.strip().lower():
        return best_match, 1.0, "Exact Copy", "Perfect match (copied)"
    elif best_score >= 0.75:
        return best_match, best_score, "Strong", "Chunked properly"
    elif best_score >= 0.5:
        return best_match, best_score, "Partial", "Partially matching"
    else:
        return best_match, best_score, "Missing", "No strong match"

def run_chunking_qc(pptx_path):
    prs = Presentation(pptx_path)
    all_rows, summary = [], []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = get_slide_text_with_position(slide)
        vo_texts = get_vo_text(slide)
        chunk_status_list, copy_match_list = [], []
        slide_points_with_order = []

        for j, point in enumerate(slide_texts, start=1):
            vo_match, score, match_type, comment = compare_point_to_vo(point, vo_texts)
            copy_match = "Yes" if match_type == "Exact Copy" else "No"
            chunk_status = "Chunked Properly" if match_type in ["Exact Copy", "Strong"] else "Not Chunked Properly" if match_type == "Missing" else "Partially Chunked"
            slide_points_with_order.append({
                "Slide Number": i,
                "Point Number": j,
                "Slide Point": point,
                "Matched VO Sentence": vo_match,
                "Similarity Score": round(score, 2),
                "Match Type": match_type,
                "Exact Copy-Paste": copy_match,
                "Comment": comment
            })
            chunk_status_list.append(chunk_status)
            copy_match_list.append(copy_match)

        if not slide_texts:
            summary.append({"Slide Number": i, "Chunking Status": "No Content", "Direct Match with Note": "No", "Comment": "No slide points found"})
        else:
            if all(c == "Chunked Properly" for c in chunk_status_list):
                chunk_status = "Chunked Properly"
            elif all(c == "Not Chunked Properly" for c in chunk_status_list):
                chunk_status = "Not Chunked Properly"
            else:
                chunk_status = "Partially Chunked"
            copy_flag = "Yes" if any(c == "Yes" for c in copy_match_list) else "No"
            summary.append({
                "Slide Number": i,
                "Chunking Status": chunk_status,
                "Direct Match with Note": copy_flag,
                "Comment": ""
            })

        for point in slide_points_with_order:
            all_rows.append(point)

    return pd.DataFrame(all_rows), pd.DataFrame(summary)
