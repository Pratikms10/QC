import re
import os
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# === CONFIG ===
VALID_FONTS = ["HelveticaNowDisplay Black", "Queens Medium", "HelveticaNowDisplay Medium", "Cambria Math", "Consolas"]
VALID_COLOURS = {"#0045C0", "#27BDBB", "#FDD900", "#FF5C00", "#117673", "#CCC1FF", "#3B4096", "#F26722", "#B9CB00", "#000000", "#FFFFFF"}

# === CLEANERS ===
def remove_instructions(text):
    pattern = re.compile(r"(?i)^.*instructions\s*to\s*gd\s*:.*$")
    lines = text.splitlines()
    cleaned_lines = []
    for line in lines:
        if pattern.match(line.strip()):
            break
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()

def clean_text_for_excel(text):
    if not isinstance(text, str):
        return text
    text = text.replace("’", "'")  # Normalize curly apostrophes
    return ''.join([c for c in text if ord(c) in range(32, 127) or ord(c) in (9, 10, 13)])

# === COLOR HELPERS ===
def rgb_to_hex(rgb):
    if not rgb:
        return "", ""
    r, g, b = rgb[0], rgb[1], rgb[2]
    hex_color = "#{:02X}{:02X}{:02X}".format(r, g, b)
    name = {
        "#000000": "Black", "#FFFFFF": "White", "#F26722": "Orange",
        "#0045C0": "Blue", "#27BDBB": "Turquoise", "#FDD900": "Yellow",
        "#117673": "Teal", "#CCC1FF": "Lavender", "#3B4096": "Indigo",
        "#B9CB00": "Lime"
    }.get(hex_color.upper(), "Custom")
    return hex_color, name

# === SHAPE EXTRACTOR ===
def extract_shape_info(shape, slide_num, file_name):
    infos = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            infos.extend(extract_shape_info(sub, slide_num, file_name))
        return infos

    font_name = font_size = text = ""
    font_hex = font_name_color = fill_hex = fill_name = line_hex = line_name = ""

    if shape.has_text_frame:
        tf = shape.text_frame
        if tf.text:
            raw_text = tf.text.strip()
            text = clean_text_for_excel(raw_text)

            try:
                # Fallback-safe loop: get first valid font name/size
                for para in tf.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font_name = font.name if font.name else "No Font"  # If no font is found, set "No Font"
                        font_size = font.size.pt if font.size else "No Size"  # If no font size, set "No Size"

                        if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
                            font_hex, font_name_color = rgb_to_hex(font.color.rgb)

                        if font_name and font_size:
                            break
                    if font_name and font_size:
                        break
            except Exception as e:
                print(f"Font extraction error in slide {slide_num}, shape {shape.name}: {e}")

    try:
        if shape.fill and shape.fill.fore_color and hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
            fill_hex, fill_name = rgb_to_hex(shape.fill.fore_color.rgb)
    except:
        pass

    try:
        if shape.line and shape.line.color and hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
            line_hex, line_name = rgb_to_hex(shape.line.color.rgb)
    except:
        pass

    infos.append({
        "File Name": file_name,
        "Slide Number": slide_num,
        "Shape Name / Table Cell": shape.name,
        "Shape Type": str(shape.shape_type),
        "Font Name": font_name,  # Ensure that the font name is correctly returned
        "Font Size": font_size,  # Ensure that the font size is properly handled
        "Font Color Hex": font_hex,
        "Font Color Name": font_name_color,
        "Fill Color Hex": fill_hex,
        "Fill Color Name": fill_name,
        "Line Color Hex": line_hex,
        "Line Color Name": line_name,
        "Extracted Text": text
    })

    return infos

# === EXTRACT NOTES ===
def extract_notes(file_path):
    prs = Presentation(file_path)
    notes = []
    for i, slide in enumerate(prs.slides, 1):
        note = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
        notes.append((i, remove_instructions(note)))
    return notes

# === WORD COMPARISON ===
def compare_words(text_a, text_b):
    a_words = re.findall(r'\b\w+\b', text_a.lower())
    b_words = re.findall(r'\b\w+\b', text_b.lower())
    return [word for word in a_words if word not in b_words]

# === SLIDE INFO EXTRACTOR ===
def extract_ppt_data(ppt_path):
    prs = Presentation(ppt_path)
    file_name = os.path.basename(ppt_path)
    all_info = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            all_info.extend(extract_shape_info(shape, idx, file_name))
    return pd.DataFrame(all_info)

# === MAIN VALIDATOR ===
def run_notes_validation(path_a, path_b):
    notes_a = extract_notes(path_a)
    notes_b = extract_notes(path_b)
    shapes_df = extract_ppt_data(path_b)
    shapes_df["Font Size"] = pd.to_numeric(shapes_df["Font Size"], errors='coerce')

    df_notes_a = pd.DataFrame([{"Slide": s, "Note Text": t} for s, t in notes_a])
    df_notes_b = pd.DataFrame([{"Slide": s, "Note Text": t} for s, t in notes_b])

    comparison_rows = []
    b_index = 0

    for slide_a, note_a in notes_a:
        matched_slides = []
        combined_b = ""
        a_words = re.findall(r'\b\w+\b', note_a.lower())
        b_words_collected = []

        for j in range(b_index, len(notes_b)):
            slide_b, note_b = notes_b[j]
            b_words = re.findall(r'\b\w+\b', note_b.lower())
            b_words_collected.extend(b_words)
            combined_b += " " + note_b
            matched_slides.append(str(slide_b))
            if all(word in b_words_collected for word in a_words):
                b_index = j + 1
                break

        missing = compare_words(note_a, combined_b)
        extra = compare_words(combined_b, note_a)
        comparison_rows.append({
            "File A Slide": slide_a,
            "File A Notes": note_a,
            "Matched B Slides": ", ".join(matched_slides),
            "Missing Words": " ".join(missing),
            "Extra Words": " ".join(extra),
            "Highlighted Notes": note_a
        })

    df_cmp = pd.DataFrame(comparison_rows)
    return df_notes_a, df_notes_b, df_cmp, shapes_df
