import win32com.client
import pythoncom
import pandas as pd
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
import re
import os

model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

def get_vo_text(slide):
    if not slide.has_notes_slide:
        return []
    notes_text = slide.notes_slide.notes_text_frame.text
    vo_section = re.search(r'VO:(.*?)(Image Link:|Instructions to GD:|$)', notes_text, re.DOTALL | re.IGNORECASE)
    if not vo_section:
        return []
    vo_clean = vo_section.group(1)
    return [line.strip("-• \n") for line in vo_clean.split("\n") if line.strip()]

# ✅ Updated to return a list of lines from text frame or grouped shapes
def extract_text_from_shape(shape):
    chunks = []
    if shape.Type == 6:  # Grouped shape
        for i in range(1, shape.GroupItems.Count + 1):
            sub_shape = shape.GroupItems(i)
            chunks.extend(extract_text_from_shape(sub_shape))
    elif shape.HasTextFrame and shape.TextFrame.HasText:
        lines = shape.TextFrame.TextRange.Text.splitlines()
        for line in lines:
            clean = line.strip("•- \n\t")
            if clean:
                chunks.append(clean)
    return chunks

# ✅ Updated to use the new chunk extraction
def get_animated_slide_points(pptx_path):
    pptx_path = os.path.abspath(pptx_path)
    pythoncom.CoInitialize()
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True

    presentation = ppt.Presentations.Open(pptx_path, WithWindow=0)
    slide_points = {}

    for slide in presentation.Slides:
        points = []
        seen_texts = set()
        for effect in slide.TimeLine.MainSequence:
            try:
                shape = effect.Shape
                chunks = extract_text_from_shape(shape)
                for chunk in chunks:
                    if chunk and chunk not in seen_texts:
                        points.append(chunk)
                        seen_texts.add(chunk)
            except Exception:
                continue
        slide_points[slide.SlideNumber] = points

    presentation.Close()
    ppt.Quit()
    pythoncom.CoUninitialize()
    return slide_points

def compare_point_to_vo(point, vo_lines, used_indices):
    if not vo_lines:
        return ("", 0.0, "Missing", "No VO content", -1)

    point_emb = model.encode(point, convert_to_tensor=True)
    vo_embs = model.encode(vo_lines, convert_to_tensor=True)
    cosine_scores = util.cos_sim(point_emb, vo_embs)[0]

    for idx in sorted(used_indices):
        cosine_scores[idx] = -1

    best_score = float(cosine_scores.max())
    best_idx = int(cosine_scores.argmax())
    best_match = vo_lines[best_idx]

    if point.strip().lower() == best_match.strip().lower():
        return best_match, 1.0, "Exact Copy", "Perfect match (copied)", best_idx
    elif best_score >= 0.75:
        return best_match, best_score, "Strong", "Chunked properly", best_idx
    elif best_score >= 0.5:
        return best_match, best_score, "Partial", "Partially matching", best_idx
    else:
        return best_match, best_score, "Missing", "No strong match", -1

def run_chunking_qc_with_animation(pptx_path):
    prs = Presentation(pptx_path)
    animated_points = get_animated_slide_points(pptx_path)
    all_rows, summary = [], []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = animated_points.get(i, [])
        vo_texts = get_vo_text(slide)
        chunk_status_list, copy_match_list = [], []
        slide_points_with_order = []
        used_vo_indices = set()
        fallback_used = False

        if not slide_texts:
            fallback_used = True
            shape_data = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    lines = shape.text.splitlines()
                    for line in lines:
                        clean = line.strip("•- \n\t")
                        if clean:
                            shape_data.append((shape.top, shape.left, clean))
            shape_data.sort(key=lambda x: (x[0], x[1]))
            slide_texts = [text for _, _, text in shape_data]

        for j, point in enumerate(slide_texts, start=1):
            vo_match, score, match_type, comment, matched_idx = compare_point_to_vo(point, vo_texts, used_vo_indices)
            if matched_idx >= 0:
                used_vo_indices.add(matched_idx)
            copy_match = "Yes" if match_type == "Exact Copy" else "No"
            chunk_status = (
                "Chunked Properly" if match_type in ["Exact Copy", "Strong"]
                else "Not Chunked Properly" if match_type == "Missing"
                else "Partially Chunked"
            )
            fallback_note = " [Fallback: No animation]" if fallback_used else ""

            slide_points_with_order.append({
                "Slide Number": i,
                "Point Number": j,
                "Slide Point": point,
                "Matched VO Sentence": vo_match,
                "Similarity Score": round(score, 2),
                "Match Type": match_type,
                "Exact Copy-Paste": copy_match,
                "Comment": comment + fallback_note
            })

            chunk_status_list.append(chunk_status)
            copy_match_list.append(copy_match)

        if not slide_texts:
            summary.append({
                "Slide Number": i,
                "Chunking Status": "No Content",
                "Direct Match with Note": "No",
                "Comment": "No slide points found"
            })
        else:
            if all(c == "Chunked Properly" for c in chunk_status_list):
                chunk_status = "Chunked Properly"
            elif all(c == "Not Chunked Properly" for c in chunk_status_list):
                chunk_status = "Not Chunked Properly"
            else:
                chunk_status = "Partially Chunked"

            copy_flag = "Yes" if any(c == "Yes" for c in copy_match_list) else "No"

            summary.append({
                "Slide Number": i,
                "Chunking Status": chunk_status,
                "Direct Match with Note": copy_flag,
                "Comment": "Fallback used" if fallback_used else ""
            })

        all_rows.extend(slide_points_with_order)

    return pd.DataFrame(all_rows), pd.DataFrame(summary)
