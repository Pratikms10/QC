import re
import pandas as pd
from pptx import Presentation

# Load US → UK dictionary once
def load_us_to_uk_dict(csv_path="us_to_uk_dictionary.csv"):
    df = pd.read_csv(csv_path)
    return {row["American"].strip().lower(): row["British"].strip().lower()
            for _, row in df.iterrows() if pd.notna(row["American"]) and pd.notna(row["British"])}

# Contraction patterns
contractions_re = re.compile(r"\b(?:[A-Za-z]+n’t|'s|'re|'ve|'ll|'d|'m)\b", re.IGNORECASE)
period_re = re.compile(r"\.")
extra_space_re = re.compile(r"\s{2,}")

def clean_text(text):
    return text.replace("\n", " ").strip()

def scan_text_issues(pptx_path, dictionary_path="us_to_uk_dictionary.csv"):
    prs = Presentation(pptx_path)
    us_to_uk = load_us_to_uk_dict(dictionary_path)
    word_re = re.compile(r'\b(' + '|'.join(re.escape(word) for word in us_to_uk.keys()) + r')\b', re.IGNORECASE)

    findings = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        entries = []

        # Extract content from slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                entries.append(("Content", shape.text.strip()))

        # Extract content from notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                entries.append(("Notes", notes))

        # Analyze all extracted text
        for loc_type, text in entries:
            text_clean = clean_text(text)

            # Check for issues
            contractions = contractions_re.findall(text_clean)
            periods = period_re.findall(text_clean)
            extra_spaces = extra_space_re.findall(text_clean)
            us_words = [m.group() for m in word_re.finditer(text_clean)]

            if contractions or periods or extra_spaces or us_words:
                findings.append({
                    "Slide Number": slide_idx,
                    "Location": loc_type,
                    "Original Text": text_clean,
                    "Contractions Used": ", ".join(set(contractions)),
                    "Ending Period Used": "." if text_clean.endswith('.') else "",
                    "Period In Middle": "Yes" if '.' in text_clean[:-1] else "",
                    "Extra Space": "Yes" if extra_spaces else "",
                    "US English Used": ", ".join(set(us_words))
                })

    return pd.DataFrame(findings)
